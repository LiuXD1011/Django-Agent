import base64
import io
import random
import shutil
import subprocess
import sys
import xml.etree.ElementTree as ET
from pathlib import Path
from types import SimpleNamespace
from unittest import skipUnless
from unittest.mock import patch
from zipfile import ZIP_DEFLATED, ZipFile

from django.test import SimpleTestCase
from PIL import Image

from personal_knowledge_base.document_parsing import parse_document
from personal_knowledge_base.document_parsing.remote_images import UnsafeRemoteImageError, validate_remote_url


FIXTURE_DIR = Path(__file__).with_name("testdata") / "legacy"
OFFICE_CONVERTER = shutil.which("soffice") or shutil.which("libreoffice")


def image_bytes(fmt="PNG", size=(96, 96)):
    image = Image.frombytes("RGB", size, random.Random(42).randbytes(size[0] * size[1] * 3))
    output = io.BytesIO()
    image.save(output, format=fmt)
    return output.getvalue()


def fixture_bytes(name):
    return (FIXTURE_DIR / name).read_bytes()


def build_xlsx_fixture():
    from openpyxl import Workbook

    workbook = Workbook()
    revenue = workbook.active
    revenue.title = "Revenue"
    revenue.append(["Quarter", "Cached", "Fallback"])
    revenue.append(["Q1", "=SUM(40,60)", "=SUM(20,30)"])
    revenue.append(["Merged note"])
    revenue.merge_cells("A3:B3")
    workbook.create_sheet("Empty")
    stream = io.BytesIO()
    workbook.save(stream)

    namespace = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
    ET.register_namespace("", namespace)
    output = io.BytesIO()
    with ZipFile(io.BytesIO(stream.getvalue())) as source, ZipFile(output, "w", ZIP_DEFLATED) as target:
        for info in source.infolist():
            payload = source.read(info.filename)
            if info.filename == "xl/worksheets/sheet1.xml":
                root = ET.fromstring(payload)
                value = root.find(f".//{{{namespace}}}c[@r='B2']/{{{namespace}}}v")
                value.text = "100"
                payload = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            target.writestr(info, payload)
    return output.getvalue()


def build_docx_fixture():
    import docx

    document = docx.Document()
    document.add_heading("Converted document", level=2)
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Revenue"
    table.cell(1, 1).text = "100"
    stream = io.BytesIO()
    document.save(stream)
    return stream.getvalue()


def build_pptx_fixture():
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "Converted presentation"
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Revenue"
    table.cell(1, 1).text = "100"
    stream = io.BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def completed_conversion(target_format, output):
    def convert(command, **kwargs):
        output_dir = Path(command[command.index("--outdir") + 1])
        input_path = Path(command[-1])
        output_path = output_dir / f"{input_path.stem}.{target_format}"
        output_path.write_bytes(output)
        return subprocess.CompletedProcess(command, 0, b"", b"")

    return convert


class DocumentParsingTests(SimpleTestCase):
    def test_xlsx_merge_ranges_are_read_from_zip_xml_relationships(self):
        from personal_knowledge_base.document_parsing.spreadsheet import _safe_xlsx_part, _xlsx_merged_ranges

        self.assertEqual(
            _xlsx_merged_ranges(build_xlsx_fixture()),
            {"Revenue": ["A3:B3"], "Empty": []},
        )
        self.assertIsNone(_safe_xlsx_part("xl/workbook.xml", "/outside.xml"))

    def test_xlsx_loads_formula_and_cached_workbooks_in_read_only_mode(self):
        import openpyxl

        with patch("openpyxl.load_workbook", wraps=openpyxl.load_workbook) as load_workbook:
            parse_document("finance.xlsx", build_xlsx_fixture())

        self.assertEqual(load_workbook.call_count, 2)
        self.assertTrue(all(call.kwargs["read_only"] for call in load_workbook.call_args_list))

    def test_xlsx_uses_cached_formula_values_with_per_cell_fallback_provenance(self):
        parsed = parse_document("finance.xlsx", build_xlsx_fixture())

        rows = [block for block in parsed.text_blocks if block.metadata.get("sheet_name") == "Revenue"]
        self.assertEqual(rows[0].metadata["headers"], ["Quarter", "Cached", "Fallback"])
        self.assertEqual(rows[1].metadata["row_start"], 2)
        self.assertEqual(rows[1].text, "Q1 | 100 | =SUM(20,30)")
        self.assertEqual(
            rows[1].metadata["cell_provenance"],
            [
                {"coordinate": "A2", "column": 1, "value": "Q1", "source": "literal"},
                {
                    "coordinate": "B2",
                    "column": 2,
                    "value": "100",
                    "source": "cached",
                    "formula": "=SUM(40,60)",
                },
                {
                    "coordinate": "C2",
                    "column": 3,
                    "value": "=SUM(20,30)",
                    "source": "formula",
                    "formula": "=SUM(20,30)",
                },
            ],
        )
        self.assertEqual(rows[1].block_type, "record")
        self.assertEqual((rows[1].source_start, rows[1].source_end), (2, 2))

    def test_xlsx_emits_merged_cell_ranges(self):
        parsed = parse_document("finance.xlsx", build_xlsx_fixture())

        revenue_rows = [block for block in parsed.text_blocks if block.metadata.get("sheet_name") == "Revenue"]
        self.assertTrue(revenue_rows)
        self.assertTrue(all(block.metadata["merged_ranges"] == ["A3:B3"] for block in revenue_rows))

    def test_xls_emits_merged_cell_ranges_and_cell_provenance(self):
        parsed = parse_document("finance.xls", fixture_bytes("sample.xls"))

        rows = [block for block in parsed.text_blocks if block.metadata.get("sheet_name") == "Revenue"]
        self.assertEqual(rows[0].metadata["headers"], ["Quarter", "Amount"])
        self.assertEqual(rows[1].metadata["row_start"], 2)
        self.assertEqual(rows[1].metadata["column_start"], 1)
        self.assertEqual(rows[1].metadata["merged_ranges"], ["A3:B3"])
        self.assertEqual(rows[1].metadata["cell_provenance"][1]["source"], "cached")
        self.assertEqual(rows[1].block_type, "record")

    @patch("personal_knowledge_base.document_parsing.legacy_office.shutil.which", return_value="/usr/bin/soffice")
    @patch("personal_knowledge_base.document_parsing.legacy_office.subprocess.run")
    def test_doc_uses_headless_conversion(self, run, _which):
        run.side_effect = completed_conversion("docx", build_docx_fixture())

        parsed = parse_document("sample.doc", fixture_bytes("sample.doc"))

        self.assertTrue(parsed.text_blocks)
        self.assertIn("--headless", run.call_args.args[0])
        self.assertEqual(parsed.text_blocks[0].text, "Converted document")
        self.assertEqual(parsed.text_blocks[0].block_type, "heading")
        self.assertEqual(parsed.text_blocks[0].metadata["heading_level"], 2)
        self.assertEqual(parsed.text_blocks[1].block_type, "table")
        self.assertEqual(parsed.text_blocks[1].metadata["rows"][1], ["Revenue", "100"])

    @patch("personal_knowledge_base.document_parsing.legacy_office.shutil.which", return_value="/usr/bin/libreoffice")
    @patch("personal_knowledge_base.document_parsing.legacy_office.subprocess.run")
    def test_ppt_uses_headless_conversion(self, run, _which):
        run.side_effect = completed_conversion("pptx", build_pptx_fixture())

        parsed = parse_document("sample.ppt", fixture_bytes("sample.ppt"))

        self.assertEqual(parsed.text_blocks[0].text, "Converted presentation")
        self.assertEqual(parsed.text_blocks[0].block_type, "text_box")
        self.assertEqual(parsed.text_blocks[0].metadata["slide_number"], 1)
        self.assertEqual(parsed.text_blocks[1].block_type, "table")
        self.assertEqual(parsed.text_blocks[1].metadata["rows"][1], ["Revenue", "100"])
        self.assertEqual(run.call_args.args[0][3], "pptx")

    @skipUnless(OFFICE_CONVERTER, "LibreOffice or soffice is required for legacy conversion integration tests")
    def test_real_doc_fixture_converts_to_a_nonempty_document(self):
        parsed = parse_document("sample.doc", fixture_bytes("sample.doc"))

        self.assertTrue(parsed.text_blocks or parsed.images)

    @skipUnless(OFFICE_CONVERTER, "LibreOffice or soffice is required for legacy conversion integration tests")
    def test_real_ppt_fixture_converts_to_a_nonempty_document(self):
        parsed = parse_document("sample.ppt", fixture_bytes("sample.ppt"))

        self.assertTrue(parsed.text_blocks or parsed.images)

    @patch("personal_knowledge_base.document_parsing.legacy_office.shutil.which", return_value=None)
    def test_legacy_converter_unavailable_has_stable_error_code(self, _which):
        with self.assertRaises(ValueError) as raised:
            parse_document("sample.doc", fixture_bytes("sample.doc"))

        self.assertEqual(raised.exception.code, "legacy_office_converter_unavailable")

    @patch("personal_knowledge_base.document_parsing.legacy_office.shutil.which", return_value="/usr/bin/soffice")
    @patch("personal_knowledge_base.document_parsing.legacy_office.subprocess.run")
    def test_legacy_conversion_timeout_has_stable_error_code(self, run, _which):
        run.side_effect = subprocess.TimeoutExpired(["soffice"], 30)

        with self.assertRaises(ValueError) as raised:
            parse_document("sample.doc", fixture_bytes("sample.doc"))

        self.assertEqual(raised.exception.code, "legacy_office_conversion_timeout")

    @patch("personal_knowledge_base.document_parsing.legacy_office.shutil.which", return_value=None)
    @patch("personal_knowledge_base.document_parsing.legacy_office.subprocess.run")
    def test_encrypted_legacy_input_is_detected_before_converter_lookup(self, run, _which):
        detector = SimpleNamespace(is_encrypted=lambda: True)
        msoffcrypto = SimpleNamespace(OfficeFile=lambda _stream: detector)

        with patch.dict(sys.modules, {"msoffcrypto": msoffcrypto}):
            with self.assertRaises(ValueError) as raised:
                parse_document("sample.doc", fixture_bytes("sample.doc"))

        self.assertEqual(raised.exception.code, "legacy_office_encrypted")
        run.assert_not_called()

    @patch("personal_knowledge_base.document_parsing.legacy_office.shutil.which", return_value="/usr/bin/soffice")
    @patch("personal_knowledge_base.document_parsing.legacy_office.subprocess.run")
    def test_protected_success_without_output_is_classified_as_encrypted(self, run, _which):
        detector = SimpleNamespace(is_encrypted=lambda: False)
        msoffcrypto = SimpleNamespace(OfficeFile=lambda _stream: detector)
        run.return_value = subprocess.CompletedProcess(
            ["soffice"],
            0,
            stdout=b"Source document is password protected",
            stderr=b"",
        )

        with patch.dict(sys.modules, {"msoffcrypto": msoffcrypto}):
            with self.assertRaises(ValueError) as raised:
                parse_document("sample.doc", fixture_bytes("sample.doc"))

        self.assertEqual(raised.exception.code, "legacy_office_encrypted")

    @patch("personal_knowledge_base.document_parsing.legacy_office.shutil.which", return_value="/usr/bin/soffice")
    @patch("personal_knowledge_base.document_parsing.legacy_office.subprocess.run")
    def test_missing_legacy_conversion_output_has_stable_error_code(self, run, _which):
        run.return_value = subprocess.CompletedProcess(["soffice"], 0, b"", b"")

        with self.assertRaises(ValueError) as raised:
            parse_document("sample.doc", fixture_bytes("sample.doc"))

        self.assertEqual(raised.exception.code, "legacy_office_conversion_output_missing")

    @patch("personal_knowledge_base.document_parsing.legacy_office.shutil.which", return_value="/usr/bin/soffice")
    @patch("personal_knowledge_base.document_parsing.legacy_office.subprocess.run")
    def test_unparseable_legacy_conversion_output_has_stable_error_code(self, run, _which):
        run.side_effect = completed_conversion("docx", b"not a docx")

        with self.assertRaises(ValueError) as raised:
            parse_document("sample.doc", fixture_bytes("sample.doc"))

        self.assertEqual(raised.exception.code, "legacy_office_converted_output_invalid")

    def test_plain_alias_uses_builtin_text_parser(self):
        parsed = parse_document("notes.txt", "第一段\n\n第二段".encode(), engine="plain")
        self.assertEqual([block.text for block in parsed.text_blocks], ["第一段\n\n第二段"])
        self.assertEqual(parsed.images, [])

    def test_standalone_image_becomes_an_image_block(self):
        parsed = parse_document("diagram.png", image_bytes())
        self.assertEqual(len(parsed.images), 1)
        self.assertEqual(parsed.images[0].source_type, "standalone")
        self.assertEqual((parsed.images[0].width, parsed.images[0].height), (96, 96))
        self.assertEqual(parsed.images[0].block_index, 0)

    def test_markdown_data_uri_preserves_order_and_extracts_image(self):
        encoded = base64.b64encode(image_bytes()).decode()
        markdown = f"前文\n\n![结构图](data:image/png;base64,{encoded})\n\n后文"
        parsed = parse_document("readme.md", markdown.encode())
        self.assertEqual(len(parsed.images), 1)
        self.assertEqual(parsed.images[0].source_type, "markdown_data")
        self.assertLess(parsed.text_blocks[0].block_index, parsed.images[0].block_index)
        self.assertLess(parsed.images[0].block_index, parsed.text_blocks[-1].block_index)

    def test_private_remote_markdown_image_is_rejected(self):
        with self.assertRaises(UnsafeRemoteImageError):
            validate_remote_url("http://127.0.0.1/private.png")

    def test_docx_emits_heading_and_preserves_inline_text_image_order(self):
        import docx

        document = docx.Document()
        document.add_heading("二级标题", level=2)
        paragraph = document.add_paragraph()
        paragraph.add_run("图片之前")
        paragraph.add_run().add_picture(io.BytesIO(image_bytes()))
        paragraph.add_run("图片之后")
        stream = io.BytesIO()
        document.save(stream)

        parsed = parse_document("manual.docx", stream.getvalue())

        ordered = parsed.ordered_blocks
        self.assertEqual(ordered[0].block_type, "heading")
        self.assertEqual(ordered[0].metadata["heading_level"], 2)
        self.assertEqual(ordered[1].text, "图片之前")
        self.assertIs(ordered[2], parsed.images[0])
        self.assertEqual(ordered[3].text, "图片之后")
        self.assertEqual(ordered[1].block_type, "paragraph")
        self.assertEqual(ordered[3].block_type, "paragraph")
        self.assertEqual(len(parsed.images), 1)
        self.assertEqual(parsed.images[0].source_type, "docx")

    def test_docx_table_emits_rows_cells_and_image_metadata(self):
        import docx

        document = docx.Document()
        table = document.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "项目"
        table.cell(0, 1).text = "金额"
        table.cell(1, 0).text = "收入"
        table.cell(1, 1).text = "100"
        table.cell(1, 1).paragraphs[0].add_run().add_picture(io.BytesIO(image_bytes()))
        stream = io.BytesIO()
        document.save(stream)

        parsed = parse_document("table.docx", stream.getvalue())

        table_block = parsed.text_blocks[0]
        self.assertEqual(table_block.block_type, "table")
        self.assertEqual(table_block.metadata["rows"], [["项目", "金额"], ["收入", "100"]])
        self.assertEqual(
            table_block.metadata["cells"],
            [
                {"row": 1, "column": 1, "text": "项目"},
                {"row": 1, "column": 2, "text": "金额"},
                {"row": 2, "column": 1, "text": "收入"},
                {"row": 2, "column": 2, "text": "100"},
            ],
        )
        self.assertEqual(len(parsed.images), 1)
        self.assertEqual(parsed.images[0].source_type, "docx")
        self.assertEqual(parsed.images[0].metadata["table_cell"], {"row": 2, "column": 2})

    def test_pptx_emits_slide_text_table_and_pictures_in_shape_order(self):
        from pptx import Presentation
        from pptx.util import Inches

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        textbox.text = "季度趋势"
        slide.shapes.add_picture(io.BytesIO(image_bytes()), Inches(1), Inches(2))
        table = slide.shapes.add_table(2, 2, Inches(3), Inches(2), Inches(4), Inches(1)).table
        table.cell(0, 0).text = "项目"
        table.cell(0, 1).text = "金额"
        table.cell(1, 0).text = "收入"
        table.cell(1, 1).text = "100"
        slide.shapes.add_picture(io.BytesIO(image_bytes()), Inches(1), Inches(4))
        second_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        second_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "第二页"
        stream = io.BytesIO()
        presentation.save(stream)

        parsed = parse_document("report.pptx", stream.getvalue())

        ordered = parsed.ordered_blocks
        self.assertEqual([block.block_index for block in ordered], list(range(5)))
        self.assertEqual(ordered[0].block_type, "text_box")
        self.assertEqual(ordered[0].metadata["slide_number"], 1)
        self.assertIs(ordered[1], parsed.images[0])
        self.assertEqual(ordered[1].metadata["shape_index"], 1)
        self.assertEqual(ordered[2].block_type, "table")
        self.assertEqual(ordered[2].metadata["rows"], [["项目", "金额"], ["收入", "100"]])
        self.assertEqual(ordered[2].metadata["cells"][3], {"row": 2, "column": 2, "text": "100"})
        self.assertIs(ordered[3], parsed.images[1])
        self.assertEqual(ordered[3].metadata["shape_index"], 3)
        self.assertEqual(ordered[4].text, "第二页")
        self.assertEqual(ordered[4].metadata["slide_number"], 2)
        self.assertTrue(all(block.metadata["slide_number"] == 1 for block in ordered[:4]))

    def test_scanned_pdf_page_is_rendered_as_image(self):
        import fitz
        pdf = fitz.open()
        page = pdf.new_page(width=300, height=300)
        page.insert_image(fitz.Rect(0, 0, 300, 300), stream=image_bytes(size=(128, 128)))
        data = pdf.tobytes()
        pdf.close()
        parsed = parse_document("scan.pdf", data)
        self.assertEqual(len(parsed.images), 1)
        self.assertEqual(parsed.images[0].source_type, "scanned_pdf")
        self.assertEqual(parsed.images[0].page_index, 0)
