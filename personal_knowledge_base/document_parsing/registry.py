import base64
import csv
import io
import json
import re
from statistics import median
from pathlib import Path

from .images import ImageTooSmallError, InvalidImageError, guess_image_mime, inspect_image
from .legacy_office import LegacyOfficeParseError, convert_legacy_office
from .remote_images import UnsafeRemoteImageError, download_remote_image
from .spreadsheet import parse_xls, parse_xlsx
from .types import ImageBlock, ParsedDocument, ParseWarning, TextBlock


IMAGE_TYPES = {"jpg", "jpeg", "png", "gif", "bmp", "tif", "tiff", "webp", "svg"}
TEXT_TYPES = {"txt", "log", "py"}
MARKDOWN_IMAGE_RE = re.compile(r"!\[([^\]]*)\]\(([^)]+)\)", re.I)
DATA_IMAGE_RE = re.compile(r"^data:(image/[^;,]+);base64,(.+)$", re.I | re.S)
PDF_NUMBERED_HEADING_RE = re.compile(r"^(?:\d+(?:\.\d+){0,5}|[IVXLC]+|[A-Z])[.)、\s]+", re.I)

_CJK_RE = re.compile(r"[\u4e00-\u9fff\u3400-\u4dbf]")
# 行尾连字符断行时，若连字符前的词干是常见复合词前缀，则连字符是词本身的一部分
# （non-static、band-pass、Human-centered），保留连字符直接拼接；其余视为排版音节断行
# （ca-tion、opti-cal），去掉连字符拼接。词表可按需扩充。
_PDF_COMPOUND_PREFIXES = frozenset({
    "non", "self", "multi", "sub", "super", "inter", "intra", "anti", "semi",
    "micro", "macro", "meta", "proto", "pseudo", "over", "under", "counter",
    "cross", "band", "phase", "human", "video", "motion", "real", "full",
    "half", "pre", "post", "out", "off", "well", "soft", "hard", "auto",
})


def _is_cjk(char: str) -> bool:
    return bool(_CJK_RE.match(char))


def _join_pdf_block_lines(lines: list) -> str:
    """把 PDF 文本块内的排版行合并为自然段落文本。

    PyMuPDF 的 dict 模式按版面把一个段落拆成多个视觉行；直接用 "\\n" 拼接会把
    行换行与断词连字符（ca-\\ntion）永久写入正文。合并规则：
    - 行尾连字符：词干是复合词前缀（见 _PDF_COMPOUND_PREFIXES）→ 保留连字符直拼；
      否则视为音节断行 → 去掉连字符直拼；
    - 相邻行以 CJK 收尾/开头 → 直拼（中文无词间空格）；
    - 其余 → 单个空格拼接。
    """
    result = ""
    for raw_line in lines:
        line = (raw_line or "").strip()
        if not line:
            continue
        if not result:
            result = line
            continue
        if result.endswith("-"):
            stem = result[:-1].rsplit(" ", 1)[-1]
            next_is_word = bool(line) and line[0].isascii() and line[0].isalpha()
            if next_is_word:
                if stem.isalpha() and stem.lower() not in _PDF_COMPOUND_PREFIXES:
                    # 音节断行（ca- tion）：去掉连字符直拼
                    result = result[:-1] + line
                else:
                    # 复合词前缀（non- static）或多连字符词干（state-of- the-art）：
                    # 连字符是词的一部分，保留并直拼
                    result = result + line
                continue
            # 连字符后不是字母（数字/符号/CJK）：保守保留连字符直拼
            result = result + line
            continue
        if _is_cjk(result[-1]) or _is_cjk(line[0]):
            result = result + line
        else:
            result = result + " " + line
    return result


def _image_block(data: bytes, mime_type: str, source_type: str, source_ref: str, block_index: int, page_index=None, metadata=None):
    width, height, detected_mime = inspect_image(data, mime_type)
    return ImageBlock(data, detected_mime, width, height, source_type, source_ref, block_index, page_index, metadata or {})


def _warning(document: ParsedDocument, code: str, message: str, block_index=None, source_ref=""):
    document.warnings.append(ParseWarning(code, message, block_index, source_ref))


def parse_text(name: str, data: bytes) -> ParsedDocument:
    suffix = Path(name).suffix.lower().lstrip(".")
    text = data.decode("utf-8", errors="ignore")
    if suffix == "json":
        try:
            text = json.dumps(json.loads(text), ensure_ascii=False, indent=2)
        except Exception:
            pass
    elif suffix == "csv":
        try:
            text = "\n".join(" | ".join(row) for row in csv.reader(io.StringIO(text)))
        except Exception:
            pass
    elif suffix in {"html", "htm"}:
        text = re.sub(r"<(script|style).*?</\1>", "", text, flags=re.I | re.S)
        text = re.sub(r"<[^>]+>", " ", text)
    return ParsedDocument(text_blocks=[TextBlock(text, 0)] if text.strip() else [])


def parse_image(name: str, data: bytes) -> ParsedDocument:
    document = ParsedDocument()
    mime_type = "image/svg+xml" if name.lower().endswith(".svg") else guess_image_mime(name)
    try:
        document.images.append(_image_block(data, mime_type, "standalone", name, 0))
    except InvalidImageError as exc:
        _warning(document, "invalid_or_small_image", str(exc), 0, name)
    return document


def parse_markdown(name: str, data: bytes) -> ParsedDocument:
    document = ParsedDocument()
    text = data.decode("utf-8", errors="ignore")
    cursor = 0
    block_index = 0
    for match in MARKDOWN_IMAGE_RE.finditer(text):
        before = text[cursor:match.start()].strip()
        if before:
            document.text_blocks.append(TextBlock(before, block_index))
            block_index += 1
        alt, target = match.group(1), match.group(2).strip().strip("<>")
        try:
            data_match = DATA_IMAGE_RE.match(target)
            if data_match:
                image_data = base64.b64decode(data_match.group(2), validate=True)
                document.images.append(_image_block(image_data, data_match.group(1), "markdown_data", alt or "data-uri", block_index))
            elif target.startswith(("http://", "https://")):
                image_data, mime_type, resolved = download_remote_image(target)
                document.images.append(_image_block(image_data, mime_type, "markdown_remote", resolved, block_index))
            else:
                document.text_blocks.append(TextBlock(match.group(0), block_index))
                _warning(document, "relative_image_unavailable", "relative Markdown image is not part of the upload", block_index, target)
        except (ValueError, InvalidImageError, UnsafeRemoteImageError, OSError) as exc:
            document.text_blocks.append(TextBlock(match.group(0), block_index))
            _warning(document, "image_fetch_failed", str(exc), block_index, target)
        block_index += 1
        cursor = match.end()
    tail = text[cursor:].strip()
    if tail:
        document.text_blocks.append(TextBlock(tail, block_index))
    return document


def _pdf_font_summary(raw_block: dict) -> tuple[list[float], list[str], bool]:
    sizes = []
    fonts = []
    bold = False
    for line in raw_block.get("lines", []):
        for span in line.get("spans", []):
            if not str(span.get("text") or "").strip():
                continue
            size = span.get("size")
            if isinstance(size, (int, float)):
                sizes.append(float(size))
            font = str(span.get("font") or "")
            if font:
                fonts.append(font)
            bold = bold or "bold" in font.lower() or bool(int(span.get("flags") or 0) & 16)
    return sizes, fonts, bold


def _pdf_repeated_margin_text(items: list[dict]) -> set[int]:
    occurrences: dict[str, set[int]] = {}
    for index, item in enumerate(items):
        if item["kind"] != "text" or len(item["text"]) > 120:
            continue
        page_height = item["page_height"]
        y0, y1 = item["bbox"][1], item["bbox"][3]
        if y1 > page_height * 0.12 and y0 < page_height * 0.88:
            continue
        normalized = " ".join(item["text"].split()).casefold()
        if normalized:
            occurrences.setdefault(normalized, set()).add(item["page_index"])
    repeated = {text for text, pages in occurrences.items() if len(pages) >= 2}
    return {
        index
        for index, item in enumerate(items)
        if item["kind"] == "text" and " ".join(item["text"].split()).casefold() in repeated
    }


def _pdf_structure(items: list[dict]) -> None:
    """Classify text blocks using only relative document typography and layout."""
    ignored = _pdf_repeated_margin_text(items)
    body_sizes = [
        median(item["font_sizes"])
        for index, item in enumerate(items)
        if item["kind"] == "text" and index not in ignored and item["font_sizes"]
    ]
    body_size = median(body_sizes) if body_sizes else 0.0
    heading_sizes = set()
    for index, item in enumerate(items):
        item["ignored"] = index in ignored
        if item["kind"] != "text" or item["ignored"]:
            continue
        size = median(item["font_sizes"]) if item["font_sizes"] else 0.0
        short = len(item["text"]) <= 160
        numbered = bool(PDF_NUMBERED_HEADING_RE.match(item["text"]))
        larger = size >= body_size + max(1.5, body_size * 0.18)
        candidate = short and (larger or (numbered and item["is_bold"] and size >= body_size))
        confidence = 0.35
        if larger:
            confidence += 0.35
        if item["is_bold"]:
            confidence += 0.15
        if numbered:
            confidence += 0.1
        if short:
            confidence += 0.05
        item["is_heading"] = candidate and confidence >= 0.7
        item["structure_confidence"] = round(min(confidence, 0.99), 2) if item["is_heading"] else 0.5
        if item["is_heading"]:
            heading_sizes.add(size)
    levels = {size: min(position + 1, 6) for position, size in enumerate(sorted(heading_sizes, reverse=True))}
    for item in items:
        if item.get("is_heading"):
            size = median(item["font_sizes"]) if item["font_sizes"] else 0.0
            item["heading_level"] = levels.get(size, 1)


def parse_pdf(name: str, data: bytes) -> ParsedDocument:
    import fitz

    document = ParsedDocument()
    pdf = fitz.open(stream=data, filetype="pdf")
    items = []
    try:
        for page_index, page in enumerate(pdf):
            page_text = page.get_text("text").strip()
            has_visual = bool(page.get_images(full=True) or page.get_drawings())
            if len(page_text) < 20 and has_visual:
                rendered = page.get_pixmap(dpi=150, alpha=False).tobytes("jpeg")
                items.append({
                    "kind": "image",
                    "value": rendered,
                    "ext": "jpeg",
                    "bbox": (0, 0, page.rect.width, page.rect.height),
                    "page_index": page_index,
                    "page_height": page.rect.height,
                    "source_type": "scanned_pdf",
                    "source_ref": f"page:{page_index + 1}",
                })
                continue
            for raw in page.get_text("dict").get("blocks", []):
                bbox = raw.get("bbox") or (0, 0, 0, 0)
                if raw.get("type") == 0:
                    value = _join_pdf_block_lines(
                        "".join(span.get("text", "") for span in line.get("spans", []))
                        for line in raw.get("lines", [])
                    ).strip()
                    if value:
                        font_sizes, font_names, is_bold = _pdf_font_summary(raw)
                        items.append({
                            "kind": "text",
                            "text": value,
                            "bbox": bbox,
                            "page_index": page_index,
                            "page_height": page.rect.height,
                            "font_sizes": font_sizes,
                            "font_names": font_names,
                            "is_bold": is_bold,
                        })
                elif raw.get("type") == 1 and raw.get("image"):
                    items.append({
                        "kind": "image",
                        "value": raw["image"],
                        "ext": raw.get("ext", "png"),
                        "bbox": bbox,
                        "page_index": page_index,
                        "page_height": page.rect.height,
                        "source_type": "pdf_embedded",
                        "source_ref": f"page:{page_index + 1}",
                    })
            for drawing_index, drawing in enumerate(page.get_drawings()):
                rect = drawing.get("rect")
                if not rect or rect.width < 64 or rect.height < 64:
                    continue
                try:
                    pixmap = page.get_pixmap(clip=rect, dpi=150, alpha=False)
                    items.append({
                        "kind": "image",
                        "value": pixmap.tobytes("png"),
                        "ext": "png",
                        "bbox": (rect.x0, rect.y0, rect.x1, rect.y1),
                        "page_index": page_index,
                        "page_height": page.rect.height,
                        "source_type": "pdf_vector",
                        "source_ref": f"page:{page_index + 1}:drawing:{drawing_index}",
                    })
                except InvalidImageError:
                    continue
    finally:
        pdf.close()
    items.sort(key=lambda item: (item["page_index"], item["bbox"][1], item["bbox"][0]))
    _pdf_structure(items)
    block_index = 0
    source_offset = 0
    has_text = False
    for item in items:
        bbox = list(item["bbox"])
        page_index = item["page_index"]
        if item["kind"] == "text":
            if item.get("ignored"):
                continue
            if has_text:
                source_offset += 2
            text = item["text"]
            metadata = {
                "bbox": bbox,
                "page_number": page_index + 1,
                "font_sizes": [round(size, 2) for size in item["font_sizes"]],
                "font_names": item["font_names"],
                "structure_confidence": item["structure_confidence"],
            }
            block_type = "heading" if item.get("is_heading") else "paragraph"
            if item.get("is_heading"):
                metadata["heading_level"] = item["heading_level"]
            document.text_blocks.append(
                TextBlock(
                    text,
                    block_index,
                    page_index,
                    metadata,
                    block_type=block_type,
                    source_start=source_offset,
                    source_end=source_offset + len(text),
                )
            )
            source_offset += len(text)
            has_text = True
        else:
            try:
                document.images.append(
                    _image_block(
                        item["value"],
                        guess_image_mime(f"image.{item['ext']}"),
                        item["source_type"],
                        item["source_ref"],
                        block_index,
                        page_index,
                        {"bbox": bbox},
                    )
                )
            except ImageTooSmallError as exc:
                _warning(document, "small_image_skipped", str(exc), block_index, item["source_ref"])
            except InvalidImageError as exc:
                _warning(document, "invalid_image", str(exc), block_index, item["source_ref"])
        block_index += 1
    return document


def parse_docx(name: str, data: bytes) -> ParsedDocument:
    import docx
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    source = docx.Document(io.BytesIO(data))
    document = ParsedDocument()
    block_index = 0

    def append_image(blip, metadata):
        nonlocal block_index
        rel_id = blip.get(qn("r:embed"))
        part = source.part.related_parts.get(rel_id)
        if not part:
            return
        try:
            document.images.append(
                _image_block(part.blob, part.content_type, "docx", rel_id, block_index, metadata=metadata)
            )
        except InvalidImageError as exc:
            _warning(document, "invalid_or_small_image", str(exc), block_index, rel_id)
        block_index += 1

    def heading_level(paragraph):
        style = paragraph.style
        style_name = getattr(style, "name", "") or ""
        style_id = getattr(style, "style_id", "") or ""
        match = re.search(r"heading\s*([1-9])$", f"{style_name} {style_id}", re.I)
        if match:
            return int(match.group(1))
        paragraph_properties = paragraph._p.pPr
        outline = paragraph_properties.find(qn("w:outlineLvl")) if paragraph_properties is not None else None
        if outline is not None:
            return int(outline.get(qn("w:val"))) + 1
        return None

    def append_paragraph(paragraph, body_index):
        nonlocal block_index
        level = heading_level(paragraph)
        style_name = getattr(paragraph.style, "name", "") or ""
        metadata = {"paragraph_index": body_index, "style_name": style_name}
        block_type = "heading" if level is not None else "paragraph"
        if level is not None:
            metadata["heading_level"] = level

        text_parts = []
        inline_position = 0

        def flush_text():
            nonlocal block_index, inline_position
            value = "".join(text_parts).strip()
            text_parts.clear()
            if not value:
                return
            fragment_metadata = {**metadata, "inline_position": inline_position}
            document.text_blocks.append(
                TextBlock(
                    value,
                    block_index,
                    metadata=fragment_metadata,
                    block_type=block_type,
                    source_start=body_index,
                    source_end=body_index,
                )
            )
            block_index += 1
            inline_position += 1

        for node in paragraph._p.iter():
            if node.tag == qn("w:t") and node.text:
                text_parts.append(node.text)
            elif node.tag == qn("w:tab"):
                text_parts.append("\t")
            elif node.tag in {qn("w:br"), qn("w:cr")}:
                text_parts.append("\n")
            elif node.tag == qn("a:blip"):
                flush_text()
                append_image(node, {**metadata, "inline_position": inline_position})
                inline_position += 1
        flush_text()

    table_index = 0
    for body_index, child in enumerate(source.element.body.iterchildren(), start=1):
        if child.tag == qn("w:p"):
            append_paragraph(Paragraph(child, source), body_index)
        elif child.tag == qn("w:tbl"):
            table = Table(child, source)
            table_index += 1
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            cells = [
                {"row": row_number, "column": column_number, "text": text}
                for row_number, row in enumerate(rows, start=1)
                for column_number, text in enumerate(row, start=1)
            ]
            metadata = {
                "table_index": table_index,
                "rows": rows,
                "cells": cells,
                "row_count": len(rows),
                "column_count": max((len(row) for row in rows), default=0),
            }
            value = "\n".join(" | ".join(row) for row in rows)
            document.text_blocks.append(
                TextBlock(
                    value,
                    block_index,
                    metadata=metadata,
                    block_type="table",
                    source_start=body_index,
                    source_end=body_index,
                )
            )
            block_index += 1

            visited_cells = set()
            for row_number, row in enumerate(table.rows, start=1):
                for column_number, cell in enumerate(row.cells, start=1):
                    if cell._tc in visited_cells:
                        continue
                    visited_cells.add(cell._tc)
                    image_metadata = {
                        "table_index": table_index,
                        "table_cell": {"row": row_number, "column": column_number},
                    }
                    for blip in cell._tc.iter(qn("a:blip")):
                        append_image(blip, image_metadata)
    return document


def parse_pptx(name: str, data: bytes) -> ParsedDocument:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    source = Presentation(io.BytesIO(data))
    document = ParsedDocument()
    block_index = 0

    def ordered_shapes(shapes, parent_path=()):
        for index, shape in enumerate(shapes):
            shape_path = (*parent_path, index)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from ordered_shapes(shape.shapes, shape_path)
            else:
                yield shape_path, shape

    for page_index, slide in enumerate(source.slides):
        slide_number = page_index + 1
        for shape_path, shape in ordered_shapes(slide.shapes):
            shape_index = shape_path[0]
            metadata = {
                "slide_number": slide_number,
                "shape_index": shape_index,
                "shape_path": list(shape_path),
                "shape_name": shape.name,
            }
            if getattr(shape, "has_table", False):
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                metadata.update(
                    {
                        "rows": rows,
                        "cells": [
                            {"row": row_number, "column": column_number, "text": text}
                            for row_number, row in enumerate(rows, start=1)
                            for column_number, text in enumerate(row, start=1)
                        ],
                        "row_count": len(rows),
                        "column_count": max((len(row) for row in rows), default=0),
                    }
                )
                value = "\n".join(" | ".join(row) for row in rows)
                document.text_blocks.append(
                    TextBlock(
                        value,
                        block_index,
                        page_index,
                        metadata,
                        block_type="table",
                        source_start=slide_number,
                        source_end=slide_number,
                    )
                )
                block_index += 1
                continue

            text = getattr(shape, "text", "").strip() if getattr(shape, "has_text_frame", False) else ""
            if text:
                document.text_blocks.append(
                    TextBlock(
                        text,
                        block_index,
                        page_index,
                        metadata,
                        block_type="text_box",
                        source_start=slide_number,
                        source_end=slide_number,
                    )
                )
                block_index += 1
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            extension = (shape.image.ext or "").lower()
            source_ref = f"slide:{slide_number}:shape:{'.'.join(str(part) for part in shape_path)}"
            if extension in {"wmf", "emf"}:
                _warning(document, "unsupported_vector_media", f".{extension} media is not supported", block_index, source_ref)
                block_index += 1
                continue
            try:
                document.images.append(
                    _image_block(
                        shape.image.blob,
                        guess_image_mime(f"image.{extension}"),
                        "pptx",
                        source_ref,
                        block_index,
                        page_index,
                        metadata,
                    )
                )
            except InvalidImageError as exc:
                _warning(document, "invalid_or_small_image", str(exc), block_index, source_ref)
            block_index += 1
    return document


def parse_legacy_office(name: str, data: bytes) -> ParsedDocument:
    converted_name, converted_data = convert_legacy_office(name, data)
    try:
        if converted_name.endswith(".docx"):
            return parse_docx(converted_name, converted_data)
        return parse_pptx(converted_name, converted_data)
    except Exception as exc:
        raise LegacyOfficeParseError(
            "legacy_office_converted_output_invalid",
            "legacy Office conversion output could not be parsed",
        ) from exc


def parse_document(name: str, data: bytes, engine: str = "builtin") -> ParsedDocument:
    if engine not in {"", "builtin", "plain"}:
        raise ValueError(f"unsupported parser engine: {engine}")
    suffix = Path(name or "").suffix.lower().lstrip(".")
    if suffix in IMAGE_TYPES:
        return parse_image(name, data)
    if suffix in {"md", "markdown"}:
        return parse_markdown(name, data)
    if suffix == "pdf":
        return parse_pdf(name, data)
    if suffix == "docx":
        return parse_docx(name, data)
    if suffix == "pptx":
        return parse_pptx(name, data)
    if suffix == "xlsx":
        return parse_xlsx(data)
    if suffix == "xls":
        return parse_xls(data)
    if suffix in {"doc", "ppt"}:
        return parse_legacy_office(name, data)
    return parse_text(name, data)
